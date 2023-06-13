import json
import logging
import os
from pptx import Presentation
import time
from Ask_gpt import send_gpt_request
import asyncio

UPLOAD_FOLDER = os.getcwd() + '\\uploads'
OUTPUT_FOLDER = os.getcwd() + '\\output'


def read_slides(presentation):
    """
    Retrieves the text content from a presentation.

    :param presentation: The ptt object from which to extract the text.
    :return: list of strings containing the combined text content from the slide's text frames.
    """
    return ['\n'.join([shape.text_frame.text for shape in slide.shapes if shape.has_text_frame])
            for slide in presentation.slides]


async def explain_slide(slide: str, slide_index: int):
    """
    get text of slide and explain it using chat-gpt

    :param slide: slide's text
    :param slide_index: the number of the slide
    :return: the response of gpt which explain slide
    """
    request = "Give me a brief explanation of this page, " \
              "Write the explanation as a short paragraph from an article:\n" + \
              "page number: " + str(slide_index) + " content: " + slide
    return await send_gpt_request(request)


async def response_handler(slides):
    """
    Handles the requests of explaining the slides to chat-gpt.

    :param slides: A list of slide contents.
    :return: A list of responses from chat-gpt for each slide.
    """
    tasks = []
    for slide_index, slide_content in enumerate(slides, start=1):
        task = asyncio.create_task(explain_slide(slide_content, slide_index))
        tasks.append(task)

    # Wait for all tasks to complete
    responses = await asyncio.gather(*tasks)
    return responses


def write_output_to_json(file_name, responses):
    """
    Writes the explanation of all slides to a JSON file.

    :param file_name: The name of the output file.
    :param responses: A list of responses from chat-gpt.
    :return: None
    """
    content_list = [response["choices"][0]["message"]["content"] for response in responses]
    slide_list = [{"slide number": i, "slide summary": content} for i, content in enumerate(content_list, start=1)]

    # write the explanation of the slides into file
    with open(file_name, "w") as file:
        json.dump(slide_list, file, indent=4)


def process_presentation(file_path):
    try:
        presentation = Presentation(file_path)
    except Exception as e:
        print("Error:", e)

    slides = read_slides(presentation)

    # get the explanation of the slides from chat gpt with async
    loop = asyncio.get_event_loop()
    responses = loop.run_until_complete(response_handler(slides))

    output_name = os.path.join(OUTPUT_FOLDER, f"{os.path.splitext(os.path.basename(file_path))[0]}.json")
    write_output_to_json(output_name, responses)


def main():
    while True:
        for file_name in os.listdir(UPLOAD_FOLDER):
            file_path = os.path.join(UPLOAD_FOLDER, file_name)  # Create the full file path
            if os.path.isfile(file_path) and file_name.lower().endswith('.pptx'):
                try:
                    logging.info("Rendering the presentation "+file_name)
                    process_presentation(file_path)
                    logging.info("Rendering succeeded")
                    os.remove(file_path)  # Delete the input file after successful processing
                except Exception as e:
                    logging.error(f"An error occurred while processing {file_path}: {str(e)}")

        time.sleep(10)


if __name__ == '__main__':
    main()
