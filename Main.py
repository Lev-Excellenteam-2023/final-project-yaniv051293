import os
from Ask_gpt import send_gpt_request
from pptx import Presentation
import asyncio
import json


def get_presentation():
    """
     Retrieves a presentation file and initializes the Presentation object.

    :return:  A dictionary containing the name of the presentation file and the Presentation object.
    """
    file_path = input("Please enter the path of the presentation file:\n")
    presentation = None
    while presentation is None:
        try:
            presentation = Presentation(file_path)
        except Exception as e:
            print("Error:", e)
            file_path = input("Please enter the correct path of the presentation file:\n")
    file_name = os.path.basename(file_path)[:-5]
    return {"name": file_name, "file": presentation}


def read_slides(presentation):
    """
    Retrieves the text content from a presentation.

    :param presentation: The ptt object from which to extract the text.
    :return: list of strings containing the combined text content from the slide's text frames.
    """
    return ['\n'.join([shape.text_frame.text for shape in slide.shapes if shape.has_text_frame])
            for slide in presentation.slides]


async def explain_slide(slide: str, slide_index: int):
    """
    get text of slide and explain it using chat-gpt

    :param slide: slide's text
    :param slide_index: the number of the slide
    :return: the response of gpt which explain slide
    """
    request = "Give me a brief explanation of this page, " \
              "Write the explanation as a short paragraph from an article:\n" + \
              "page number: " + str(slide_index) + " content: " + slide
    return await send_gpt_request(request)


async def response_handler(slides):
    """
    Handles the requests of explaining the slides to chat-gpt.

    :param slides: A list of slide contents.
    :return: A list of responses from chat-gpt for each slide.
    """
    tasks = []
    for slide_index, slide_content in enumerate(slides, start=1):
        task = asyncio.create_task(explain_slide(slide_content, slide_index))
        tasks.append(task)

    # Wait for all tasks to complete
    responses = await asyncio.gather(*tasks)
    return responses


def write_output_to_json(file_name, responses):
    """
    Writes the explanation of all slides to a JSON file.

    :param file_name: The name of the output file.
    :param responses: A list of responses from chat-gpt.
    :return: None
    """
    content_list = [response["choices"][0]["message"]["content"] for response in responses]
    slide_list = [{"slide number": i, "slide content": content} for i, content in enumerate(content_list, start=1)]

    # write the explanation of the slides into file
    with open(file_name + ".json", "w") as file:
        json.dump(slide_list, file, indent=4)


def main():
    presentation = get_presentation()
    slides = read_slides(presentation["file"])

    # get the explanation of the slides from chat gpt with async
    loop = asyncio.get_event_loop()
    responses = loop.run_until_complete(response_handler(slides))

    write_output_to_json(presentation["name"], responses)


if __name__ == '__main__':
    main()
